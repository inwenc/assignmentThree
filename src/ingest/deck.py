"""Asynchronous slide-deck ingestion for PDF and PPTX sources.

Each vector retains its one-based ``slide`` number, which is the durable
citation locator for a deck.  The flow deliberately follows the paper/video
flows: source acquisition and parsing happen in Prefect tasks, embedding uses
the already-configured text embedding provider, and deterministic Qdrant IDs
make retries replace partial work rather than duplicate it.
"""
from __future__ import annotations

import urllib.parse
import urllib.request
from dataclasses import dataclass
from pathlib import Path

from prefect import flow, task

from .. import db, storage
from ..config import TEXT_EMBED_VERSION
from ..rag import vector_store
from ..rag.embeddings import embed_docs
from .fetch import scratch_dir, sha256_file


@dataclass(frozen=True)
class Slide:
    """Text extracted from a one-based slide number."""

    number: int
    text: str


def _normalise(text: str) -> str:
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


def extract_pdf_slides(path: Path) -> list[Slide]:
    """Treat each PDF page as one slide while preserving page/slide order."""
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise RuntimeError("PDF deck ingestion requires PyMuPDF (pip install pymupdf)") from exc

    slides: list[Slide] = []
    with fitz.open(path) as pdf:
        for index, page in enumerate(pdf):
            text = _normalise(page.get_text("text", sort=True))
            if text:
                slides.append(Slide(index + 1, text))
    return slides


def extract_pptx_slides(path: Path) -> list[Slide]:
    """Extract visible text, tables, and speaker notes from a PPTX deck."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PPTX ingestion requires python-pptx (pip install python-pptx)") from exc

    presentation = Presentation(path)
    slides: list[Slide] = []
    for index, slide in enumerate(presentation.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                parts.append(shape.text)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        # Notes often hold the presenter-authored explanation for sparse slides.
        notes = getattr(slide, "notes_slide", None)
        note_frame = getattr(notes, "notes_text_frame", None) if notes else None
        if note_frame and note_frame.text.strip():
            parts.append(note_frame.text)
        text = _normalise("\n".join(parts))
        if text:
            slides.append(Slide(index + 1, text))
    return slides


def extract_slides(path: Path) -> list[Slide]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_slides(path)
    if suffix == ".pptx":
        return extract_pptx_slides(path)
    raise ValueError("deck must be a PDF or PPTX file")


def chunk_slides(slides: list[Slide], *, max_chars: int = 1_600,
                 overlap_chars: int = 180) -> list[dict]:
    """Produce one or more chunks per slide without mixing slide locators."""
    if max_chars <= 0 or not 0 <= overlap_chars < max_chars:
        raise ValueError("max_chars must be positive and overlap_chars must be smaller")
    chunks: list[dict] = []
    for slide in slides:
        text = " ".join(slide.text.split())
        start = 0
        while start < len(text):
            end = min(len(text), start + max_chars)
            if end < len(text):
                boundary = text.rfind(" ", start, end)
                if boundary > start:
                    end = boundary
            piece = text[start:end].strip()
            if piece:
                chunks.append({"slide": slide.number, "text": piece})
            if end >= len(text):
                break
            start = max(end - overlap_chars, start + 1)
    return chunks


def _deck_suffix(storage_key: str | None, url: str | None) -> str:
    source = storage_key or (urllib.parse.urlparse(url or "").path)
    suffix = Path(source).suffix.lower()
    if suffix not in {".pdf", ".pptx"}:
        raise ValueError("deck URI must end in .pdf or .pptx")
    return suffix


def fetch_deck(storage_key: str | None, url: str | None, deck_id: str) -> Path:
    """Download a stored or public PDF/PPTX deck into worker scratch space."""
    suffix = _deck_suffix(storage_key, url)
    destination = scratch_dir() / f"{deck_id}{suffix}"
    if storage_key:
        return storage.download_to(storage_key, destination)
    if not url or not url.lower().startswith(("https://", "http://")):
        raise ValueError("deck requires a storage_key or an http(s) URL")
    request = urllib.request.Request(url, headers={"User-Agent": "MomentSearch/1.0"})
    try:
        with urllib.request.urlopen(request, timeout=60) as response, destination.open("wb") as out:
            while data := response.read(1 << 20):
                out.write(data)
    except Exception:
        destination.unlink(missing_ok=True)
        raise
    return destination


@task(name="fetch-deck", retries=2, retry_delay_seconds=[30, 120])
def t_fetch_deck(deck_id: str, user_id: str) -> str:
    row = db.get_video(deck_id)
    if row is None:
        raise ValueError(f"no manifest row for {deck_id}")
    db.set_status(deck_id, "parsing", progress=0.0)
    path = fetch_deck(row.get("storage_key"), row.get("url"), deck_id)
    source_hash = sha256_file(path)
    db.set_status(deck_id, "parsing", source_hash=source_hash, progress=0.05)
    duplicate = db.find_duplicate(user_id, source_hash, exclude_id=deck_id)
    if duplicate:
        path.unlink(missing_ok=True)
        db.set_status(deck_id, "skipped", error=f"duplicate of {duplicate['id']}")
        return ""
    return str(path)


@task(name="parse-deck", retries=1, retry_delay_seconds=30)
def t_parse_deck(deck_id: str, path: str) -> list[Slide]:
    db.set_status(deck_id, "parsing", progress=0.1)
    slides = extract_slides(Path(path))
    if not slides:
        raise RuntimeError("No selectable text could be extracted from this deck.")
    db.set_progress(deck_id, 1.0)
    return slides


@task(name="chunk-deck")
def t_chunk_deck(deck_id: str, slides: list[Slide]) -> list[dict]:
    db.set_status(deck_id, "chunking", progress=0.0)
    chunks = chunk_slides(slides)
    if not chunks:
        raise RuntimeError("Deck extraction produced no indexable chunks.")
    db.set_progress(deck_id, 1.0)
    return chunks


@task(name="enrich-deck")
def t_enrich_deck(deck_id: str, chunks: list[dict]) -> list[dict]:
    """Lifecycle seam for optional image-caption/LLM enrichment.

    The existing application has no standalone captioning interface, so this
    preserves extracted text and speaker notes today.  A vision-caption stage
    can be added here without changing slide parsing, locators, or indexing.
    """
    db.set_status(deck_id, "enriching", progress=1.0)
    return chunks


@task(name="embed-index-deck", retries=2, retry_delay_seconds=60)
def t_embed_index_deck(deck_id: str, user_id: str, chunks: list[dict]) -> int:
    db.set_status(deck_id, "embedding", progress=0.0)
    vector_store.ensure_text_collection()
    vector_store.delete_video(user_id, deck_id)
    batch_size = 128
    for start in range(0, len(chunks), batch_size):
        batch = chunks[start:start + batch_size]
        vectors = embed_docs([chunk["text"] for chunk in batch])
        vector_store.upsert_chunks(user_id, deck_id, vectors, payloads=[
            {"user_id": user_id, "video_id": deck_id, "source_id": deck_id,
             "kind": "deck", "modality": "text", "slide": chunk["slide"],
             "text": chunk["text"], "embed_version": TEXT_EMBED_VERSION}
            for chunk in batch
        ])
        db.set_progress(deck_id, min(1.0, (start + len(batch)) / len(chunks)))
    db.set_status(deck_id, "indexed", frame_count=len(chunks),
                  embed_version=TEXT_EMBED_VERSION, progress=1.0)
    return len(chunks)


@flow(name="ms-ingest-deck", log_prints=True, timeout_seconds=3600)
def ingest_deck(deck_id: str, user_id: str) -> dict:
    """Prefect flow for an already-registered deck manifest row."""
    attempt = db.bump_attempts(deck_id)
    path: str | None = None
    try:
        path = t_fetch_deck(deck_id, user_id)
        if not path:
            return {"deck_id": deck_id, "skipped": True}
        slides = t_parse_deck(deck_id, path)
        chunks = t_chunk_deck(deck_id, slides)
        chunks = t_enrich_deck(deck_id, chunks)
        count = t_embed_index_deck(deck_id, user_id, chunks)
        print(f"[deck] {deck_id} indexed: {count} chunks ({len(slides)} slides, attempt {attempt})")
        return {"deck_id": deck_id, "slides": len(slides), "chunks": count}
    except Exception as exc:
        db.set_status(deck_id, "failed", error=f"{type(exc).__name__}: {exc}")
        raise
    finally:
        if path:
            Path(path).unlink(missing_ok=True)
